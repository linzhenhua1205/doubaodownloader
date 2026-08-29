#!/usr/bin/env python3
"""web-ppt-builder: HTML → PPTX 通用转换器（深蓝科技风，16:9）
基于 web-ppt-builder 模板渲染规则，覆盖全部组件：
cover / part过渡页 / table / flow(steps) / grid(cards) / points(li) / imgcard / impact四维 / src

用法:
    python3 build_pptx.py <index.html> <输出.pptx> [图片目录]

- 图片目录缺省时自动下载：按 HTML 中 <img> 出现顺序 curl 下载到
  tmp/ppt-pptx-images/<输出文件名>/img1.jpg ... imgN.jpg（需 curl，失败自动跳过不嵌图）
- 依赖: python3 -m pip install python-pptx beautifulsoup4 lxml
- QA: 用 python-pptx 读取验证页数/图片数/关键文本；无 LibreOffice 环境时不做像素级预览
"""
import re, html as htmlmod, os, sys, subprocess, tempfile
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 参数 ----------
if len(sys.argv) < 3:
    print(__doc__)
    sys.exit(1)
SRC = sys.argv[1]
OUT = sys.argv[2]
IMG_DIR = sys.argv[3] if len(sys.argv) > 3 else os.path.join(
    'tmp', 'ppt-pptx-images', os.path.splitext(os.path.basename(OUT))[0])
os.makedirs(IMG_DIR, exist_ok=True)

# 深蓝科技风配色
C_BG       = RGBColor(0x0D, 0x1B, 0x2A)   # 页面底色
C_BG2      = RGBColor(0x12, 0x29, 0x4A)   # 卡片底色
C_PANEL    = RGBColor(0x16, 0x2E, 0x52)   # 面板/表头
C_BORDER   = RGBColor(0x2C, 0x4A, 0x6E)
C_TITLE    = RGBColor(0xFF, 0xFF, 0xFF)
C_BODY     = RGBColor(0xD8, 0xE2, 0xF0)
C_MUTED    = RGBColor(0x88, 0x99, 0xAA)
C_ACCENT   = RGBColor(0x4F, 0xC3, 0xF7)   # 青色强调
C_RED      = RGBColor(0xFF, 0x52, 0x52)
C_GREEN    = RGBColor(0x69, 0xF0, 0xAE)
C_GOLD     = RGBColor(0xFF, 0xD7, 0x4F)

FONT = 'Microsoft YaHei'
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]

# ---------- 工具函数 ----------
def add_slide():
    return prs.slides.add_slide(BLANK)

def set_bg(slide, color=C_BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()
    r.shadow.inherit = False
    return r

def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    return tf

def set_run(run, text, size, color=C_BODY, bold=False, italic=False, font=FONT):
    run.text = text
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    f.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', font)

def add_para(tf, first=False, align=PP_ALIGN.LEFT, space_before=0, space_after=0):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before); p.space_after = Pt(space_after)
    return p

def rich_text(tf, html_frag, size, base_color=C_BODY, bold=False, first=False,
              align=PP_ALIGN.LEFT, space_before=0, space_after=0, line=None):
    """解析内联 HTML(b/span.red/green/hl/a) 追加到段落"""
    p = add_para(tf, first=first, align=align, space_before=space_before, space_after=space_after)
    if line: p.line_spacing = line
    # 用 bs4 解析片段
    soup = BeautifulSoup(f'<span>{html_frag}</span>', 'html.parser')
    def walk(node, cur_bold=bold, cur_color=base_color):
        for child in node.children:
            if getattr(child, 'name', None) is None:
                txt = str(child)
                if txt:
                    r = p.add_run(); set_run(r, txt, size, cur_color, cur_bold)
                continue
            name = child.name
            nb, nc = cur_bold, cur_color
            if name == 'b': nb = True
            elif name in ('strong',): nb = True
            elif name == 'span':
                cls = ' '.join(child.get('class', []))
                if 'red' in cls: nc = C_RED
                elif 'green' in cls: nc = C_GREEN
                elif 'hl' in cls: nc = C_ACCENT
                elif 'hl2' in cls: nc = C_MUTED
                elif 'tag' in cls: nc = C_GOLD
            elif name == 'a': nc = C_ACCENT
            elif name == 'br':
                r = p.add_run(); set_run(r, '\n', size, cur_color, cur_bold)
                continue
            walk(child, nb, nc)
    walk(soup)
    return p

def fit_font(text, width_in, size, bold=False, min_size=8):
    """估算字号：中文字符≈size宽，ASCII≈size*0.55，粗体≈*1.05"""
    w = 0.0
    for ch in text:
        w += size * 1.0 if ord(ch) > 0x2E80 else size * 0.55
    if bold: w *= 1.05
    available = width_in * 72.0
    if w <= available: return size
    return max(min_size, int(size * available / w))

def add_rect(slide, x, y, w, h, fill=C_BG2, line=C_BORDER, radius=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp

def add_arrow(slide, x, y, w, h, color=C_BORDER):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def strip_html(s):
    return htmlmod.unescape(re.sub(r'<[^>]+>', '', s))

def img_local(src):
    """URL → 本地图片路径（IMG_MAP: url→完整本地路径）"""
    return IMG_MAP.get(src, '')

# ---------- 页面渲染 ----------
def render_cover(slide, sec):
    set_bg(slide)
    kicker = sec.select_one('.kicker')
    title = sec.select_one('.title')
    meta = sec.select_one('.meta')
    src = sec.select_one('.src')
    if kicker:
        tf = tb(slide, Inches(0.9), Inches(0.75), Inches(11.5), Inches(0.4))
        rich_text(tf, str(kicker), 13, C_ACCENT, first=True)
    if title:
        txt = strip_html(str(title))
        # 拆分行：<br> 分隔
        parts = re.split(r'<br\s*/?>', str(title))
        y = Inches(1.5)
        for i, part in enumerate(parts):
            clean = strip_html(part).strip()
            if not clean: continue
            sz = 34 if i == 0 else (20 if 'style' in part and '26px' in part else 26)
            if 'hl2' in part: sz = 17; col = C_MUTED
            elif 'hl' in part: sz = 30; col = C_ACCENT
            else: col = C_TITLE
            tf = tb(slide, Inches(0.9), y, Inches(11.5), Inches(0.9))
            rich_text(tf, part, sz, col, bold=True, first=True, line=1.15)
            y += Inches(0.85 if i == 0 else 0.6)
    if meta:
        tf = tb(slide, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.3))
        rich_text(tf, str(meta), 13, C_BODY, first=True, line=1.4)
    if src:
        tf = tb(slide, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5))
        rich_text(tf, str(src), 10, C_MUTED, first=True)

def render_part(slide, sec):
    """过渡页：big 大标题 + sub 引导问题"""
    set_bg(slide, C_BG)
    kicker = sec.select_one('.kicker')
    big = sec.select_one('.big')
    sub = sec.select_one('.sub')
    if kicker:
        tf = tb(slide, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.4))
        rich_text(tf, str(kicker), 13, C_ACCENT, bold=True, first=True)
    if big:
        txt = strip_html(str(big))
        sz = fit_font(txt, 11.5, 26, bold=True)
        tf = tb(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.0))
        rich_text(tf, str(big), sz, C_TITLE, bold=True, first=True, line=1.25)
    if sub:
        tf = tb(slide, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.2))
        rich_text(tf, str(sub), 16, C_BODY, first=True, line=1.4)

def render_table(slide, table_el, x, y, w, h):
    rows = table_el.find_all('tr')
    nrows = len(rows)
    if nrows == 0: return
    ncols = max(len(r.find_all(['td', 'th'])) for r in rows)
    # 字号自适应
    total_chars = sum(len(strip_html(str(c))) for r in rows for c in r.find_all(['td', 'th']))
    base = 11 if total_chars < 500 else (9.5 if total_chars < 1200 else 8.5)
    row_h = h / nrows
    gfx = slide.shapes.add_table(nrows, ncols, x, y, w, h)
    tbl = gfx.table
    # 列宽
    widths = [w / ncols] * ncols
    for i, cw in enumerate(widths):
        tbl.columns[i].width = int(cw)
    # 先关闭默认样式
    tbl.first_row = False; tbl.horz_banding = False
    for ri, tr in enumerate(rows):
        cells = tr.find_all(['td', 'th'])
        tbl.rows[ri].height = int(row_h)
        for ci in range(ncols):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Pt(4)
            cell.margin_top = cell.margin_bottom = Pt(1)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tfc = cell.text_frame; tfc.word_wrap = True
            tfc.clear()
            if ri == 0 or tr.find('th'):
                fill = C_PANEL
                col = C_TITLE; bold = True
            else:
                fill = C_BG2
                col = C_BODY; bold = False
            cell.fill.solid(); cell.fill.fore_color.rgb = fill
            txt = str(cells[ci]) if ci < len(cells) else ''
            # 行首红色标记（P0等）
            rich_text(tfc, txt, base, col, bold=bold, first=True, line=1.05)
    return gfx

def render_flow(slide, flow_el, x, y, w, h):
    """水平流程: step → arrow → step"""
    steps = flow_el.select('.step')
    arrows = flow_el.select('.arrow')
    n = len(steps) + len(arrows)
    if n == 0: return
    step_w = (w - 0.35 * len(arrows)) / max(len(steps), 1)
    cx = x
    for i, st in enumerate(steps):
        # step 块
        rect = add_rect(slide, cx, y, step_w, h, fill=C_BG2, line=C_BORDER, radius=True)
        tfe = rect.text_frame; tfe.word_wrap = True
        tfe.margin_left = tfe.margin_right = Pt(6)
        tfe.margin_top = tfe.margin_bottom = Pt(4)
        st_t = st.select_one('.st'); nm = st.select_one('.nm'); ds = st.select_one('.ds')
        if st_t:
            rich_text(tfe, str(st_t), 12, C_GOLD, bold=True, first=True)
        if nm:
            rich_text(tfe, str(nm), 11, C_TITLE, bold=True, space_before=2, line=1.05)
        if ds:
            rich_text(tfe, str(ds), 9.5, C_MUTED, space_before=2, line=1.05)
        cx += step_w
        if i < len(arrows):
            aw = Inches(0.35)
            add_arrow(slide, cx, y + h/2 - Inches(0.12), aw, Inches(0.24))
            cx += aw

def render_cards(slide, grid_el, x, y, w, h):
    """卡片网格 grid g2/g3"""
    cards = grid_el.select('.card')
    if not cards: return
    g = grid_el.get('class', [])
    ncol = 3 if 'g3' in g else 2
    n = len(cards)
    nrow = (n + ncol - 1) // ncol
    cw = (w - 0.25 * (ncol - 1)) / ncol
    ch = (h - 0.2 * (nrow - 1)) / nrow
    for i, card in enumerate(cards):
        r_, c_ = i // ncol, i % ncol
        cx = x + c_ * (cw + 0.25); cy = y + r_ * (ch + 0.2)
        hot = 'hot' in card.get('class', []) or 'opp' in card.get('class', [])
        fill = C_RED if hot else C_BG2
        rect = add_rect(slide, cx, cy, cw, ch, fill=fill if not hot else RGBColor(0x3A, 0x1A, 0x24),
                        line=C_BORDER, radius=True)
        tfe = rect.text_frame; tfe.word_wrap = True
        tfe.margin_left = tfe.margin_right = Pt(8)
        tfe.margin_top = tfe.margin_bottom = Pt(6)
        h4 = card.select_one('h4')
        if h4:
            rich_text(tfe, str(h4), 12.5, C_TITLE, bold=True, first=True)
        for li in card.select('li'):
            rich_text(tfe, str(li), 10, C_BODY, space_before=3, line=1.1)

def render_points(slide, points_el, x, y, w, h):
    tfe = tb(slide, x, y, w, h)
    lis = points_el.select('li')
    if not lis: return
    total = sum(len(strip_html(str(li))) for li in lis)
    base = 12 if total < 300 else (10.5 if total < 700 else 9.5)
    for li in lis:
        rich_text(tfe, '▸ ' + str(li), base, C_BODY, first=(li is lis[0]), space_after=4, line=1.15)

def render_imgcard(slide, imgcard_el, x, y, w, h):
    img = imgcard_el.select_one('img')
    cap = imgcard_el.select_one('.cap')
    src_url = img.get('src', '') if img else ''
    local = img_local(src_url)
    img_h = Inches(1.35)
    if local and os.path.exists(local):
        pic = slide.shapes.add_picture(local, x, y, width=min(Inches(3.4), w))
        img_h = pic.height
        # 高度限制
        if pic.height > Inches(1.9):
            ratio = Inches(1.9) / pic.height
            pic.height = Inches(1.9); pic.width = int(pic.width * ratio)
            img_h = Inches(1.9)
    cy = y + img_h + Inches(0.08)
    if cap:
        tfe = tb(slide, x, cy, w, h - img_h - Inches(0.08))
        rich_text(tfe, str(cap), 8.5, C_MUTED, first=True, line=1.1)

def render_impact(slide, impact_el, x, y, w, h):
    """四维影响：sum 全宽 + 3 列"""
    dims = impact_el.select('.dim')
    if not dims: return
    # sum 单独一行
    sum_dim = impact_el.select_one('.dim.sum')
    others = [d for d in dims if d != sum_dim]
    if sum_dim:
        tfe = tb(slide, x, y, w, Inches(1.9))
        h4 = sum_dim.select_one('h4')
        if h4:
            rich_text(tfe, str(h4), 14, C_GOLD, bold=True, first=True, space_after=3)
        for li in sum_dim.select('li'):
            rich_text(tfe, str(li), 11, C_BODY, space_after=2, line=1.12)
        y += Inches(2.1)
    if others:
        n = len(others)
        cw = (w - 0.25 * (n - 1)) / n
        for i, dim in enumerate(others):
            cx = x + i * (cw + 0.25)
            rect = add_rect(slide, cx, y, cw, h - (y - Inches(1.55)), fill=C_BG2, line=C_BORDER, radius=True)
            tfe = rect.text_frame; tfe.word_wrap = True
            tfe.margin_left = tfe.margin_right = Pt(8)
            tfe.margin_top = tfe.margin_bottom = Pt(6)
            h4 = dim.select_one('h4')
            if h4:
                rich_text(tfe, str(h4), 12, C_ACCENT, bold=True, first=True, space_after=3)
            lis = dim.select('li')
            total = sum(len(strip_html(str(li))) for li in lis)
            base = 10 if total < 400 else 9
            for li in lis:
                rich_text(tfe, '• ' + str(li), base, C_BODY, space_after=2.5, line=1.12)

def render_blocks(slide, blocks, x, y, w, h, src_h=Inches(0.28)):
    """通用块渲染器：按顺序渲染任意块列表，返回结束 y 位置"""
    by = y
    for blk in blocks:
        cls = blk.get('class', [])
        if 'flow' in cls:
            render_flow(slide, blk, x, by, w, Inches(1.5))
            by += Inches(1.7)
        elif blk.name == 'table':
            rows = len(blk.find_all('tr'))
            th = max(Inches(0.9), Inches(min(rows * 0.42, 4.2)))
            render_table(slide, blk, x, by, w, th)
            by += th + Inches(0.15)
        elif 'grid' in cls:
            # grid 内嵌：可能是 cards / 混合(div>table, div>points, imgcard...)
            cards = blk.select(':scope > .card')
            if cards:
                ncol = 3 if 'g3' in cls else 2
                nrow = (len(cards) + ncol - 1) // ncol
                gh = Inches(min(nrow * 1.15, 3.4))
                render_cards(slide, blk, x, by, w, gh)
                by += gh + Inches(0.15)
            else:
                # 混合布局：按列渲染每个直接子 div（整列作为块，保留 imgcard 类）
                cols = [c for c in blk.children if getattr(c, 'name', None) == 'div' and 'src' not in c.get('class', [])]
                if cols:
                    ncol = len(cols)
                    cw = (w - 0.25 * (ncol - 1)) / ncol
                    col_h = Inches(4.6)
                    for i, col in enumerate(cols):
                        cx = x + i * (cw + 0.25)
                        # 整列作为块递归（imgcard/table 容器/points 容器均可识别）
                        render_blocks(slide, [col], cx, by, cw, col_h, src_h=Inches(0))
                    by += col_h + Inches(0.15)
        elif 'points' in cls:
            total = sum(len(strip_html(str(li))) for li in blk.select('li'))
            ph = Inches(min(0.5 + total / 260.0, 2.6))
            render_points(slide, blk, x, by, w, ph)
            by += ph + Inches(0.12)
        elif 'imgcard' in cls:
            render_imgcard(slide, blk, x, by, w, Inches(1.9))
            by += Inches(2.0)
        elif blk.name == 'div':
            # 嵌套 div：递归渲染其子块（处理 imgcard 等包裹结构）
            sub = [c for c in blk.children if getattr(c, 'name', None) in ('table', 'div', 'p', 'ul', 'ol', 'h4')]
            if sub:
                by2 = render_blocks(slide, sub, x, by, w, h - (by - y), src_h=Inches(0))
                by = max(by, by2)
        elif 'impact' in cls:
            render_impact(slide, blk, x, by, w, Inches(6.6))
            by += Inches(6.4)
        elif blk.name in ('ul', 'ol'):
            total = sum(len(strip_html(str(li))) for li in blk.find_all('li'))
            ph = Inches(min(0.5 + total / 260.0, 2.6))
            render_points(slide, blk, x, by, w, ph)
            by += ph + Inches(0.12)
        elif blk.name == 'p':
            tfe = tb(slide, x, by, w, Inches(0.5))
            rich_text(tfe, str(blk), 11, C_BODY, first=True, line=1.15)
            by += Inches(0.5)
        elif blk.name == 'h2':
            # 目录页多个 h2
            tfe = tb(slide, x + Inches(0.1), by, w, Inches(0.55))
            rich_text(tfe, str(blk), 15, C_TITLE, bold=True, first=True, line=1.1)
            by += Inches(0.55)
    return by

def render_slide(slide, sec):
    """内容页：kicker + h2 + lead + 内容块 + src"""
    set_bg(slide)
    kicker = sec.select_one('.kicker')
    h2s = sec.select('h2')
    lead = sec.select_one('.lead')
    src = sec.select_one('.src')
    is_agenda = len(h2s) > 1  # 目录页
    y = Inches(0.28)
    if kicker:
        tfe = tb(slide, Inches(0.5), y, Inches(12.3), Inches(0.3))
        rich_text(tfe, str(kicker), 10.5, C_ACCENT, bold=True, first=True)
        y += Inches(0.32)
    if h2s and not is_agenda:
        h2 = h2s[0]
        txt = strip_html(str(h2))
        sz = fit_font(txt, 12.3, 21, bold=True)
        tfe = tb(slide, Inches(0.5), y, Inches(12.3), Inches(0.75))
        rich_text(tfe, str(h2), sz, C_TITLE, bold=True, first=True, line=1.1)
        y += Inches(0.62)
    if lead:
        txt = strip_html(str(lead))
        tfe = tb(slide, Inches(0.5), y, Inches(12.3), Inches(0.6))
        rich_text(tfe, str(lead), 11.5, C_BODY, first=True, line=1.15)
        y += Inches(0.5)
    # 内容区
    content_y = y
    # 收集块
    blocks = []
    for child in sec.children:
        if getattr(child, 'name', None) is None: continue
        if child.name == 'span': continue
        if child.name == 'p' and 'lead' in child.get('class', []): continue
        if child.name == 'div' and 'src' in child.get('class', []): continue
        if child.name == 'div' and 'kicker' in child.get('class', []): continue
        if is_agenda and child.name == 'h2':
            blocks.append(child); continue
        if not is_agenda and child.name == 'h2': continue
        if child.name == 'div' and 'line' in child.get('class', []): continue
        blocks.append(child)
    src_h = Inches(0.3) if src else Inches(0)
    render_blocks(slide, blocks, Inches(0.5), content_y, Inches(12.3), Inches(7.0) - (content_y - Inches(0.2)), src_h=src_h)
    if src:
        tfe = tb(slide, Inches(0.5), Inches(7.12), Inches(12.3), Inches(0.3))
        rich_text(tfe, str(src), 8.5, C_MUTED, first=True)

# ---------- 主流程 ----------
soup = BeautifulSoup(open(SRC, encoding='utf-8').read(), 'html.parser')
sections = soup.select('section.slide, section.cover, section.part')
print(f"共 {len(sections)} 页")

IMG_MAP = {}  # url → local
# 从 HTML 提取图片顺序并映射；本地缺失时自动 curl 下载（失败跳过，render 端有 exists 兜底）
urls = [img.get('src') for img in soup.select('img')]
img_counter = 1
for u in urls:
    local = os.path.join(IMG_DIR, f'img{img_counter}.jpg')
    if not os.path.exists(local) and u:
        try:
            subprocess.run(
                ['curl', '-sL', '--max-time', '20', '-A', 'Mozilla/5.0',
                 '-o', local, u], check=False)
        except FileNotFoundError:
            print("⚠ curl 不可用，跳过图片下载")
    IMG_MAP[u] = local
    img_counter += 1

for i, sec in enumerate(sections):
    cls = sec.get('class', [])
    slide = add_slide()
    if 'cover' in cls:
        render_cover(slide, sec)
    elif 'part' in cls:
        render_part(slide, sec)
    else:
        render_slide(slide, sec)
    print(f"  [{i+1}/{len(sections)}] {sec.get('data-title', '')}")

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"\n✅ 已保存: {OUT} ({os.path.getsize(OUT)//1024} KB)")
